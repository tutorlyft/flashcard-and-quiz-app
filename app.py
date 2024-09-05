import streamlit as st
import pandas as pd
from io import StringIO
import docx
import pptx
import PyPDF2
from openai import OpenAI
import os
import logging
from PIL import Image
import pytesseract
import io
from youtube_transcript_api import YouTubeTranscriptApi
from pytube import YouTube
import sys
import traceback
import json
#from dotenv import load_dotenv

# Load environment variables from .env file
#load_dotenv()

# Initialize the OpenAI client with the API key from the environment variable
#client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])

# Set up logging
log_stream = io.StringIO()
logging.basicConfig(level=logging.INFO, stream=log_stream, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)


def extract_text(uploaded_file):
    file_extension = uploaded_file.name.split('.')[-1].lower()
    
    if file_extension in ['csv', 'xls', 'xlsx']:
        if file_extension == 'csv':
            df = pd.read_csv(uploaded_file)
        else:
            df = pd.read_excel(uploaded_file)
        return df.to_string()
    
    elif file_extension in ['txt', 'rtf']:
        return uploaded_file.getvalue().decode("utf-8")
    
    elif file_extension == 'pdf':
        try:
            pdf_reader = PyPDF2.PdfReader(uploaded_file)
            text = "\n".join([page.extract_text() for page in pdf_reader.pages])
            logger.info(f"Extracted {len(text)} characters from PDF")
            if not text.strip():
                logger.warning("Extracted text from PDF is empty")
            return text
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {str(e)}")
            logger.error(traceback.format_exc())
            return f"Error processing PDF: {str(e)}"
    
    elif file_extension in ['doc', 'docx']:
        doc = docx.Document(uploaded_file)
        return "\n".join([para.text for para in doc.paragraphs])
    
    elif file_extension in ['ppt', 'pptx']:
        prs = pptx.Presentation(uploaded_file)
        return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text')])
    
    elif file_extension in ['jpg', 'jpeg', 'png']:
        try:
            image = Image.open(io.BytesIO(uploaded_file.read()))
            text = pytesseract.image_to_string(image)
            if text.strip():
                return text
            else:
                return "No text could be extracted from the image. The image might not contain readable text, or the text might be too complex for the OCR to recognize."
        except Exception as e:
            return f"Error processing image: {str(e)}"
    
    else:
        return "Unsupported file type"

def extract_text_from_video(video_url):
    try:
        # Extract video ID from URL
        video = YouTube(video_url)
        video_id = video.video_id
        
        logger.info(f"Extracting transcript for video ID: {video_id}")
        
        # Get the transcript
        transcript = YouTubeTranscriptApi.get_transcript(video_id)
        
        # Combine all text from the transcript
        full_text = ' '.join([entry['text'] for entry in transcript])
        
        if not full_text.strip():
            logger.warning("Extracted transcript is empty.")
            return "Error: Extracted transcript is empty."
        
        logger.info(f"Successfully extracted transcript. Length: {len(full_text)} characters")
        return full_text
    except Exception as e:
        error_message = f"Error extracting text from video: {str(e)}"
        logger.error(error_message)
        return error_message

def generate_flashcards(text, num_cards=5):
    logger.info(f"Starting generate_flashcards function with {num_cards} cards")
    logger.info(f"Input text length: {len(text)}")
    
    if not client.api_key:
        logger.error("OpenAI API key is not set")
        return []

    # Split the generation into batches
    batch_size = 10
    num_batches = (num_cards + batch_size - 1) // batch_size

    all_flashcards = []

    for batch in range(num_batches):
        cards_in_batch = min(batch_size, num_cards - batch * batch_size)
        max_tokens = min(4000, cards_in_batch * 200)  # Adjust token limit per batch

        prompt = f"""
        Create {cards_in_batch} high-quality flashcards from the following text. For each flashcard:
        1. Identify a key concept, fact, or idea from the text.
        2. Create a concise question that tests understanding of this concept.
        3. Provide a clear and informative answer.
        4. Ensure the question and answer are directly related to the main ideas in the text.
        5. Vary the types of questions (e.g., definitions, explanations, comparisons, applications).
        6. Use LaTeX notation for any mathematical expressions, enclosed in $ symbols for inline math or $$ for display math.

        Format each flashcard as:
        Q: [question]
        A: [answer]

        Text: {text[:4000]}  # Limit text to 4000 characters per batch

        Flashcards:
        """

        try:
            logger.info(f"Sending request to OpenAI API for batch {batch+1} with {cards_in_batch} flashcards")
            response = client.chat.completions.create(
                model="gpt-4",
                messages=[
                    {"role": "system", "content": "You are an expert educator creating insightful, challenging flashcards to test and reinforce key concepts."},
                    {"role": "user", "content": prompt}
                ],
                max_tokens=max_tokens,
                n=1,
                temperature=0.7,
            )
            logger.info(f"Received response from OpenAI API for batch {batch+1}")
            logger.info(f"Full API response: {response}")  # Log the full response

            flashcards_text = response.choices[0].message.content.strip()
            logger.info(f"Generated flashcards text for batch {batch+1}: {flashcards_text}")  # Log the full flashcards text

            for card in flashcards_text.split("\n\n"):
                if card.startswith("Q: ") and "A: " in card:
                    question, answer = card.split("A: ")
                    question = question.replace("Q: ", "").strip()
                    answer = answer.strip()
                    all_flashcards.append((question, answer))

        except Exception as e:
            logger.error(f"Error in generate_flashcards batch {batch+1}: {str(e)}")
            logger.error(traceback.format_exc())

    logger.info(f"Total flashcards generated: {len(all_flashcards)}")
    return all_flashcards[:num_cards]

def generate_quiz(text, num_questions=5):
    logger.info(f"Starting generate_quiz function with {num_questions} questions")
    logger.info(f"Input text length: {len(text)}")
    
    if not text.strip():
        logger.error("Input text is empty")
        return [], "Input text is empty"

    if not client.api_key:
        logger.error("OpenAI API key is not set")
        return [], "OpenAI API key is not set"

    prompt = f"""
    Create a multiple-choice quiz with exactly {num_questions} questions based on the following text. It is crucial that you generate exactly {num_questions} questions. No more, no less.
    For each question:
    1. Identify a key concept or fact from the text.
    2. Create a clear and concise question about this concept.
    3. Provide four answer choices, including the correct answer and three plausible distractors.
    4. Indicate the correct answer.
    5. Use LaTeX notation for any mathematical expressions, enclosed in $ symbols for inline math.

    Format each question as:
    Q: [question]
    A: [option 1]
    B: [option 2]
    C: [option 3]
    D: [option 4]
    Correct: [letter of correct answer]

    Ensure that each option is on a separate line and starts with the corresponding letter (A, B, C, or D).
    For any mathematical expressions, use LaTeX notation. For example, x^2 should be written as $x^2$.

    Text: {text[:4000]}  # Limit text to 4000 characters

    Quiz:
    """

    try:
        logger.info("Sending request to OpenAI API")
        response = client.chat.completions.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": "You are an expert educator creating challenging multiple-choice questions to test understanding of key concepts."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=2000,
            n=1,
            temperature=0.7,
        )
        logger.info("Received response from OpenAI API")
        logger.info(f"Full API response: {response}")  # Log the full response

        if not response.choices:
            logger.error("No choices in the response")
            return [], "No choices in the OpenAI API response"

        quiz_text = response.choices[0].message.content.strip()
        logger.info(f"Generated quiz text: {quiz_text}")  # Log the full quiz text
        
        quiz = parse_quiz_text(quiz_text)
        
        logger.info(f"Parsed {len(quiz)} questions")
        
        if len(quiz) != num_questions:
            logger.warning(f"Generated {len(quiz)} questions instead of the requested {num_questions}")
        
        if not quiz:
            logger.error("No questions could be parsed from the OpenAI API response")
            return [], "No questions could be parsed from the OpenAI API response"
        
        return quiz, None
    except Exception as e:
        logger.error(f"Error in generate_quiz: {str(e)}")
        logger.error(traceback.format_exc())
        return [], f"Error in generate_quiz: {str(e)}"

def parse_quiz_text(quiz_text):
    quiz = []
    current_question = None
    current_options = []
    current_correct = None

    for line in quiz_text.split('\n'):
        line = line.strip()
        if line.startswith('Q:'):
            if current_question:
                if len(current_options) == 4 and current_correct:
                    quiz.append((current_question, current_options, current_correct))
                else:
                    logger.warning(f"Skipping malformed question: {current_question}")
            current_question = line[2:].strip()
            current_options = []
            current_correct = None
        elif line.startswith(('A:', 'B:', 'C:', 'D:')):
            current_options.append(line[2:].strip())
        elif line.startswith('Correct:'):
            correct_letter = line[8:].strip()
            if correct_letter in 'ABCD':
                current_correct = current_options[ord(correct_letter) - ord('A')]
            else:
                logger.warning(f"Invalid correct answer format: {line}")

    # Add the last question
    if current_question and len(current_options) == 4 and current_correct:
        quiz.append((current_question, current_options, current_correct))

    return quiz

import re

def clean_latex(text):
    # Remove duplicate LaTeX expressions (rendered and unrendered versions)
    pattern = r'(\$[^$]+\$)(\1)'
    cleaned = re.sub(pattern, r'\1', text)
    
    return cleaned

def main():
    st.set_page_config(page_title="AI Study Tool", page_icon="📚", layout="centered")

    st.markdown("""
    <script src="https://cdnjs.cloudflare.com/ajax/libs/mathjax/2.7.5/MathJax.js?config=TeX-MML-AM_CHTML"></script>
    <script>
    MathJax.Hub.Config({
        tex2jax: {
            inlineMath: [['$', '$'], ['\\(', '\\)']],
            displayMath: [['$$', '$$'], ['\\[', '\\]']],
            processEscapes: true,
            processEnvironments: true,
        },
        TeX: {
            extensions: ["AMSmath.js", "AMSsymbols.js"]
        },
        CommonHTML: { linebreaks: { automatic: true } },
        "HTML-CSS": { linebreaks: { automatic: true } },
        SVG: { linebreaks: { automatic: true } }
    });

    document.addEventListener("DOMContentLoaded", function() {
        setInterval(function() {
            MathJax.Hub.Queue(["Typeset", MathJax.Hub]);
        }, 1000);
    });
    </script>
    """, unsafe_allow_html=True)

    # Add custom CSS (Your existing CSS code here)
    st.markdown("""
    <style>
    /* Your existing CSS code */
    </style>
    """, unsafe_allow_html=True)

    st.title("📚 AI Study Tool")
    st.write("Welcome to the AI Study Tool! Follow the steps below to create your study materials.")
    
    def render_content(content):
        import re

        def format_latex(match):
            latex = match.group(1)
            latex = re.sub(r'([a-zA-Z])(\d+)', r'\1^{\2}', latex)
            latex = re.sub(r'(?<=[0-9a-zA-Z}])([-+])', r' \1 ', latex)
            return f'$${latex}$$'

        formatted_content = re.sub(r'\$(.+?)\$', format_latex, content)
        final_content = re.sub(r'(\$\$.*?\$\$)', r'\1', formatted_content)
        
        return st.markdown(final_content, unsafe_allow_html=True)

    # Initialize session state
    if 'step' not in st.session_state:
        st.session_state.step = 1
        st.session_state.input_method = "File Upload"
        st.session_state.extracted_text = ""
        st.session_state.study_tool = None
        st.session_state.num_items = None
        st.session_state.study_material = None
        st.session_state.quiz_answers = {}
        st.session_state.quiz_checked = {}
        st.session_state.current_card_index = 0

    # Step 1: Choose input method and provide input
    if st.session_state.step == 1:
        st.header("Step 1: Choose Your Input Method and Provide Input")
        
        tab1, tab2, tab3 = st.tabs(["File Upload", "Video Link", "Text Input"])
        
        with tab1:
            st.session_state.input_method = "File Upload"
            uploaded_file = st.file_uploader("Choose a file", type=["pdf", "doc", "docx", "ppt", "pptx", "xls", "xlsx", "txt", "csv", "rtf", "jpg", "jpeg", "png"])
            if uploaded_file is not None:
                st.session_state.extracted_text = extract_text(uploaded_file)
                st.success(f"File '{uploaded_file.name}' uploaded and processed successfully!")
        
        with tab2:
            st.session_state.input_method = "Video Link"
            video_url = st.text_input("Enter YouTube video URL:")
            if video_url:
                st.session_state.extracted_text = extract_text_from_video(video_url)
                if st.session_state.extracted_text.startswith("Error:"):
                    st.error(st.session_state.extracted_text)
                else:
                    st.success("Video processed successfully!")
        
        with tab3:
            st.session_state.input_method = "Text Input"
            text_input = st.text_area("Enter your text here:")
            if text_input:
                st.session_state.extracted_text = text_input
                st.success("Text input received successfully!")

        if st.button("Next", key="step1_next"):
            if st.session_state.extracted_text:
                st.session_state.step = 2
                st.rerun()
            else:
                st.error("Please provide input before proceeding.")

    # Step 2: Choose study tool
    elif st.session_state.step == 2:
        st.header("Step 2: Choose Your Study Tool")
        st.session_state.study_tool = st.radio("Select a study tool:", ["Flashcards", "Quiz"], key="study_tool_radio")
        
        col1, col2 = st.columns(2)
        with col1:
            if st.button("Back", key="step2_back"):
                st.session_state.step = 1
                st.rerun()
        with col2:
            if st.button("Next", key="step2_next"):
                st.session_state.step = 3
                st.rerun()

    # Step 3: Select number of items
    elif st.session_state.step == 3:
        st.header(f"Step 3: Select Number of {st.session_state.study_tool}")
        if st.session_state.study_tool == "Flashcards":
            st.session_state.num_items = st.selectbox("Select number of flashcards to generate:", [5, 10, 25, 50], key="flashcard_num")
        else:
            st.session_state.num_items = st.selectbox("Select number of quiz questions to generate:", [5, 10, 15, 20, 25], key="quiz_num")
        
        col1, col2 = st.columns(2)
        with col1:
            if st.button("Back", key="step3_back"):
                st.session_state.step = 2
                st.rerun()
        with col2:
            if st.button(f"Create My {st.session_state.study_tool}", key="create_material"):
                st.session_state.step = 4
                st.session_state.study_material = None  # Reset study material when creating new
                st.rerun()

    # Step 4: Generate and display study material
    elif st.session_state.step == 4:
        st.header(f"Your {st.session_state.study_tool}")

        # Check if extracted text is present
        if not st.session_state.extracted_text:
            st.error("No input text provided. Please provide some text and try again.")
            return
        
        # Only generate study material if it hasn't been generated yet
        if st.session_state.study_material is None:
            with st.spinner(f"Generating your {st.session_state.study_tool.lower()}..."):
                try:
                    logger.info(f"Attempting to generate {st.session_state.study_tool}")
                    logger.info(f"Extracted text length: {len(st.session_state.extracted_text)}")

                    if st.session_state.study_tool == "Flashcards":
                        st.session_state.study_material = generate_flashcards(st.session_state.extracted_text, st.session_state.num_items)
                    else:  # Quiz
                        st.session_state.study_material, error_message = generate_quiz(st.session_state.extracted_text, st.session_state.num_items)
                    
                    if st.session_state.study_material:
                        logger.info(f"Generated study material: {st.session_state.study_material}")
                        st.success(f"Successfully generated {st.session_state.study_tool}!")
                    else:
                        logger.warning("No study material generated from input.")
                        st.warning("No study material generated. Please try again.")
                        if error_message:
                            st.error(f"Error: {error_message}")
                except Exception as e:
                    logger.error(f"Error generating study material: {str(e)}")
                    logger.error(traceback.format_exc())
                    st.error(f"An error occurred while generating {st.session_state.study_tool.lower()}. Please try again.")

        if st.session_state.study_material:
            logger.info(f"Displaying {st.session_state.study_tool}")
            if st.session_state.study_tool == "Flashcards":
                st.write(f"Generated {len(st.session_state.study_material)} Flashcards:")
                
                if 'current_card_index' not in st.session_state:
                    st.session_state.current_card_index = 0

                flashcard_container = st.container()

                with flashcard_container:
                    if 0 <= st.session_state.current_card_index < len(st.session_state.study_material):
                        current_question, current_answer = st.session_state.study_material[st.session_state.current_card_index]
                        
                        col1, col2, col3 = st.columns([1, 3, 1])
                        with col2:
                            st.markdown("### Question")
                            st.markdown(clean_latex(current_question), unsafe_allow_html=True)
                            if st.button("Reveal Answer"):
                                st.markdown("### Answer")
                                st.markdown(clean_latex(current_answer), unsafe_allow_html=True)

                # Navigation buttons
                col1, col2, col3 = st.columns(3)
                with col1:
                    if st.button("Previous", disabled=st.session_state.current_card_index == 0):
                        st.session_state.current_card_index = max(0, st.session_state.current_card_index - 1)
                        st.rerun()
                with col3:
                    if st.button("Next", disabled=st.session_state.current_card_index == len(st.session_state.study_material) - 1):
                        st.session_state.current_card_index = min(len(st.session_state.study_material) - 1, st.session_state.current_card_index + 1)
                        st.rerun()

                # Progress indicator
                st.progress((st.session_state.current_card_index + 1) / len(st.session_state.study_material))
                st.write(f"Card {st.session_state.current_card_index + 1} of {len(st.session_state.study_material)}")
                            
            else:  # Quiz
                st.write(f"Generated Quiz with {len(st.session_state.study_material)} Questions:")
                for i, (question, options, correct) in enumerate(st.session_state.study_material, 1):
                    st.subheader(f"Question {i}")
                    st.markdown(clean_latex(question), unsafe_allow_html=True)
                    
                    answer_key = f"q_{i}"
                    if answer_key not in st.session_state.quiz_answers:
                        st.session_state.quiz_answers[answer_key] = None

                    # Use radio buttons for answer selection
                    user_answer = st.radio(
                        f"Select your answer for Question {i}:",
                        options,
                        key=f"radio_{answer_key}",
                        index=None,
                        format_func=lambda x: clean_latex(x)
                    )
                    
                    # Update the answer immediately when selected
                    if user_answer is not None:
                        st.session_state.quiz_answers[answer_key] = user_answer
                    
                    check_button = st.button(f"Check Answer", key=f"check_{i}", disabled=user_answer is None)
                    
                    if check_button and user_answer is not None:
                        st.session_state.quiz_checked[answer_key] = True
                        if st.session_state.quiz_answers[answer_key] == correct:
                            st.success("Correct! 🎉")
                        else:
                            st.error("Incorrect.")
                            st.markdown(f"The correct answer is: **{clean_latex(correct)}**", unsafe_allow_html=True)
                    
                    st.write("---")

                # Calculate and display the score after all questions are answered
                answered_questions = sum(1 for key in st.session_state.quiz_checked if st.session_state.quiz_checked[key])
                if answered_questions == len(st.session_state.study_material):
                    correct_answers = sum(1 for i, (_, _, correct) in enumerate(st.session_state.study_material, 1)
                                        if st.session_state.quiz_answers[f"q_{i}"] == correct)
                    score = (correct_answers / len(st.session_state.study_material)) * 100
                    st.success(f"Quiz completed! Your score: {score:.2f}%")

        if st.button("Start Over", key="start_over"):
            logger.info("Starting over")
            for key in list(st.session_state.keys()):
                del st.session_state[key]
            st.session_state.step = 1
            st.rerun()

if __name__ == "__main__":
    main()